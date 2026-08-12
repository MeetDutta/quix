import os
import json
import math
import hashlib
from typing import List, Dict, Any, Optional
import pypdf
import docx
import pptx
import google.generativeai as genai
from app.config import settings

class RAGService:
    def __init__(self):
        self.vector_store_path = os.path.join(settings.UPLOAD_DIR, "vector_store.json")
        self.vectors: List[Dict[str, Any]] = []
        self._load_vectors()
        
        self.api_key = settings.GEMINI_API_KEY
        self.enabled = bool(self.api_key)
        if self.enabled:
            genai.configure(api_key=self.api_key)

    def _load_vectors(self):
        """Loads vectors from local JSON file."""
        if os.path.exists(self.vector_store_path):
            try:
                with open(self.vector_store_path, "r") as f:
                    self.vectors = json.load(f)
            except Exception:
                self.vectors = []

    def _save_vectors(self):
        """Saves vectors to local JSON file."""
        os.makedirs(os.path.dirname(self.vector_store_path), exist_ok=True)
        with open(self.vector_store_path, "w", encoding="utf-8") as f:
            json.dump(self.vectors, f, ensure_ascii=False)

    def _sanitize_unicode(self, text: Optional[str]) -> str:
        """
        Cleans lone surrogate code points (\ud800-\udfff) and unencodable characters
        extracted from complex PDFs or non-standard fonts so they encode cleanly to UTF-8.
        """
        if not text or not isinstance(text, str):
            return ""
        try:
            # Replaces lone surrogate code points with clean character
            cleaned = text.encode("utf-16", "surrogatepass").decode("utf-16", "replace")
            # Strip remaining replacement characters if any, and ensure valid UTF-8
            return cleaned.encode("utf-8", "ignore").decode("utf-8", "ignore")
        except Exception:
            return "".join(c for c in text if not (0xD800 <= ord(c) <= 0xDFFF))

    def extract_text(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """
        Extracts text from PDF, DOCX, PPTX, or Images.
        For images or complex pages, falls back to Gemini multi-modal text extraction.
        Sanitizes all extracted strings against Unicode surrogates.
        """
        ext = os.path.splitext(filename.lower())[1]
        pages = []

        try:
            if ext == ".pdf":
                reader = pypdf.PdfReader(file_path)
                for page_idx, page in enumerate(reader.pages):
                    try:
                        raw_text = page.extract_text() or ""
                    except Exception:
                        raw_text = ""
                    clean_text = self._sanitize_unicode(raw_text)
                    if clean_text.strip():
                        pages.append({"page_number": page_idx + 1, "text": clean_text})
            elif ext == ".docx":
                doc = docx.Document(file_path)
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                raw_text = "\n".join(full_text)
                clean_text = self._sanitize_unicode(raw_text)
                if clean_text.strip():
                    pages.append({"page_number": 1, "text": clean_text})
            elif ext == ".pptx":
                prs = pptx.Presentation(file_path)
                for slide_idx, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    raw_text = "\n".join(slide_text)
                    clean_text = self._sanitize_unicode(raw_text)
                    if clean_text.strip():
                        pages.append({"page_number": slide_idx + 1, "text": clean_text})
            elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
                # Use Gemini Vision model directly for OCR to bypass tesseract binary limits!
                if self.enabled:
                    model = genai.GenerativeModel("models/gemini-3.5-flash")
                    # Read binary image data
                    with open(file_path, "rb") as img_file:
                        img_data = img_file.read()
                    image_part = {
                        "mime_type": f"image/{ext[1:] if ext[1:] != 'jpg' else 'jpeg'}",
                        "data": img_data
                    }
                    response = model.generate_content([
                        image_part, 
                        "Extract all readable text, titles, numbers, diagrams description, and facts from this image exactly."
                    ])
                    raw_text = response.text
                else:
                    raw_text = f"[Mock OCR Text for {filename}]"
                clean_text = self._sanitize_unicode(raw_text)
                pages.append({"page_number": 1, "text": clean_text})
            elif ext in [".txt", ".csv", ".md"]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as txt_file:
                    raw_text = txt_file.read()
                clean_text = self._sanitize_unicode(raw_text)
                if clean_text.strip():
                    pages.append({"page_number": 1, "text": clean_text})
        except Exception as e:
            raise ValueError(f"Error parsing file {filename}: {str(e)}")

        return pages

    def chunk_text(self, pages: List[Dict[str, Any]], chunk_size: int = 800, overlap: int = 150) -> List[Dict[str, Any]]:
        """
        Splits extracted pages into overlapping chunks for semantic retrieval.
        """
        chunks = []
        chunk_idx = 0
        
        for p in pages:
            text = self._sanitize_unicode(p["text"])
            page_num = p["page_number"]
            
            # Simple word-based chunker
            words = text.split()
            i = 0
            while i < len(words):
                chunk_words = words[i:i + chunk_size]
                chunk_text = " ".join(chunk_words)
                
                chunks.append({
                    "chunk_index": chunk_idx,
                    "page_number": page_num,
                    "content": self._sanitize_unicode(chunk_text)
                })
                
                chunk_idx += 1
                i += (chunk_size - overlap)
                
        return chunks

    def compute_embedding(self, text: str) -> List[float]:
        """
        Computes text embedding vector using Google Gemini model models/text-embedding-004.
        Falls back to a hashed float array if AI is not enabled.
        """
        sanitized_text = self._sanitize_unicode(text)
        if self.enabled:
            try:
                result = genai.embed_content(
                    model="models/text-embedding-004",
                    content=sanitized_text,
                    task_type="retrieval_document"
                )
                return result["embedding"]
            except Exception:
                pass
                
        # Simple deterministic fallback vector of dimension 768
        h = hashlib.sha256(sanitized_text.encode("utf-8")).digest()
        vector = []
        for index in range(768):
            val = h[index % len(h)] / 255.0
            vector.append(val)
        return vector

    def add_document_to_index(self, doc_id: str, doc_title: str, chunks: List[Dict[str, Any]], subject_id: Optional[str] = None):
        """
        Generates embeddings for chunks and appends them to local index.
        Stores subject_id for strict subject-level vector search scoping.
        """
        for c in chunks:
            clean_content = self._sanitize_unicode(c["content"])
            vector = self.compute_embedding(clean_content)
            self.vectors.append({
                "chunk_id": f"{doc_id}_{c['chunk_index']}",
                "document_id": doc_id,
                "subject_id": subject_id,
                "doc_title": self._sanitize_unicode(doc_title),
                "content": clean_content,
                "page_number": c["page_number"],
                "embedding": vector
            })
        self._save_vectors()

    def remove_document_from_index(self, doc_id: str):
        """
        Removes all chunks of a specific document from index.
        """
        self.vectors = [v for v in self.vectors if v["document_id"] != doc_id]
        self._save_vectors()

    def search_similarity(
        self, 
        query: str, 
        limit: int = 5, 
        document_ids: Optional[List[str]] = None,
        subject_id: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Finds the top similarity matches using Cosine Similarity.
        Filterable by dynamic list of document IDs or specific subject_id.
        """
        query_vector = self.compute_embedding(query)
        scored_chunks = []
        
        # Filter vectors if document_ids or subject_id are provided
        target_vectors = self.vectors
        if subject_id:
            target_vectors = [v for v in target_vectors if v.get("subject_id") == subject_id]
        if document_ids:
            target_vectors = [v for v in target_vectors if v.get("document_id") in document_ids]
            
        for v in target_vectors:
            score = self._cosine_similarity(query_vector, v["embedding"])
            scored_chunks.append({
                "chunk_id": v["chunk_id"],
                "content": v["content"],
                "page_number": v["page_number"],
                "document_id": v["document_id"],
                "subject_id": v.get("subject_id"),
                "doc_title": v["doc_title"],
                "score": score
            })
            
        # Sort by score descending
        scored_chunks.sort(key=lambda x: x["score"], reverse=True)
        return scored_chunks[:limit]

    def _cosine_similarity(self, vec1: List[float], vec2: List[float]) -> float:
        """Helper to calculate cosine similarity between two float vectors."""
        dot_product = sum(a * b for a, b in zip(vec1, vec2))
        norm_a = math.sqrt(sum(a * a for a in vec1))
        norm_b = math.sqrt(sum(b * b for b in vec2))
        if norm_a == 0 or norm_b == 0:
            return 0.0
        return dot_product / (norm_a * norm_b)
